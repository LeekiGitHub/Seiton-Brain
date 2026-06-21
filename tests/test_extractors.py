from pathlib import Path
from unittest.mock import MagicMock, patch

from docx import Document
from pptx import Presentation
from pypdf import PdfWriter

from app.vault.extractors import (
    PDF_NO_TEXT_TYPE,
    SUPPORTED_EXTENSIONS,
    DocxExtractor,
    MarkdownExtractor,
    PdfExtractor,
    PlainTextExtractor,
    PptxExtractor,
    get_extractor,
    is_supported,
)


def test_markdown_extractor_parses_frontmatter(tmp_path: Path):
    note = tmp_path / "Idea.md"
    note.write_text(
        "---\ntitle: Fitness App\ncategory: idea\n---\n\n# Heading\n\nTrack workouts.",
        encoding="utf-8",
    )
    doc = MarkdownExtractor().extract(note)
    assert doc.title == "Fitness App"
    assert doc.category == "idea"
    assert doc.doc_type == "markdown"
    assert "Track workouts." in doc.text
    assert "title: Fitness App" not in doc.text  # frontmatter entfernt


def test_markdown_extractor_falls_back_to_stem(tmp_path: Path):
    note = tmp_path / "No Frontmatter.md"
    note.write_text("Just body text.", encoding="utf-8")
    doc = MarkdownExtractor().extract(note)
    assert doc.title == "No Frontmatter"
    assert doc.category == ""
    assert doc.text == "Just body text."


def test_plaintext_extractor(tmp_path: Path):
    note = tmp_path / "Rechnung 2026.txt"
    note.write_text("Betrag: 42 EUR", encoding="utf-8")
    doc = PlainTextExtractor().extract(note)
    assert doc.title == "Rechnung 2026"
    assert doc.text == "Betrag: 42 EUR"
    assert doc.doc_type == "text"
    assert doc.category == ""


def test_get_extractor_resolves_by_suffix(tmp_path: Path):
    assert isinstance(get_extractor(tmp_path / "a.md"), MarkdownExtractor)
    assert isinstance(get_extractor(tmp_path / "a.MARKDOWN"), MarkdownExtractor)
    assert isinstance(get_extractor(tmp_path / "a.txt"), PlainTextExtractor)
    assert isinstance(get_extractor(tmp_path / "scan.pdf"), PdfExtractor)
    assert isinstance(get_extractor(tmp_path / "doc.PDF"), PdfExtractor)
    assert isinstance(get_extractor(tmp_path / "Bewerbung.docx"), DocxExtractor)
    assert isinstance(get_extractor(tmp_path / "Bewerbung.DOCX"), DocxExtractor)
    assert isinstance(get_extractor(tmp_path / "Pitch.pptx"), PptxExtractor)
    assert get_extractor(tmp_path / "photo.jpg") is None
    assert get_extractor(tmp_path / "alt.doc") is None  # altes Binaerformat nicht unterstuetzt


def test_is_supported_and_extensions():
    assert is_supported(Path("a.md"))
    assert is_supported(Path("a.txt"))
    assert is_supported(Path("a.pdf"))
    assert is_supported(Path("a.docx"))
    assert is_supported(Path("a.pptx"))
    assert not is_supported(Path("a.jpg"))
    assert not is_supported(Path("a.doc"))
    assert ".md" in SUPPORTED_EXTENSIONS
    assert ".txt" in SUPPORTED_EXTENSIONS
    assert ".pdf" in SUPPORTED_EXTENSIONS
    assert ".docx" in SUPPORTED_EXTENSIONS
    assert ".pptx" in SUPPORTED_EXTENSIONS
    assert ".jpg" not in SUPPORTED_EXTENSIONS


def _fake_pdf_reader(pages_text: list[str], title: str | None = None) -> MagicMock:
    reader = MagicMock()
    reader.metadata = MagicMock(title=title)
    reader.pages = [MagicMock(extract_text=MagicMock(return_value=t)) for t in pages_text]
    return reader


def test_pdf_extractor_with_text_layer(tmp_path: Path):
    pdf = tmp_path / "Bewerbung.pdf"
    pdf.write_bytes(b"%PDF-1.4")
    with patch(
        "app.vault.extractors.PdfReader",
        return_value=_fake_pdf_reader(["Seite eins.", "Seite zwei."], title="Anschreiben"),
    ):
        doc = PdfExtractor().extract(pdf)
    assert doc.title == "Anschreiben"  # aus PDF-Metadaten
    assert doc.doc_type == "pdf"
    assert "Seite eins." in doc.text
    assert "Seite zwei." in doc.text


def test_pdf_extractor_title_falls_back_to_stem(tmp_path: Path):
    pdf = tmp_path / "Zeugnis 2024.pdf"
    pdf.write_bytes(b"%PDF-1.4")
    with patch(
        "app.vault.extractors.PdfReader",
        return_value=_fake_pdf_reader(["Inhalt"], title=None),
    ):
        doc = PdfExtractor().extract(pdf)
    assert doc.title == "Zeugnis 2024"


def test_pdf_extractor_no_text_layer_marks_for_ocr(tmp_path: Path):
    """Echter Scan ohne Text-Layer (leere Seite) → pdf_no_text."""
    pdf = tmp_path / "Scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    with pdf.open("wb") as fh:
        writer.write(fh)

    doc = PdfExtractor().extract(pdf)
    assert doc.doc_type == PDF_NO_TEXT_TYPE
    assert doc.text == ""
    assert doc.title == "Scan"


def test_pdf_extractor_corrupt_file_does_not_raise(tmp_path: Path):
    pdf = tmp_path / "kaputt.pdf"
    pdf.write_bytes(b"not really a pdf")
    doc = PdfExtractor().extract(pdf)
    assert doc.doc_type == PDF_NO_TEXT_TYPE
    assert doc.text == ""


def test_docx_extractor_paragraphs_and_tables(tmp_path: Path):
    path = tmp_path / "Bewerbung.docx"
    document = Document()
    document.core_properties.title = "Anschreiben"
    document.add_paragraph("Sehr geehrte Damen und Herren,")
    document.add_paragraph("ich bewerbe mich auf die Stelle.")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Position"
    table.rows[0].cells[1].text = "Entwickler"
    document.save(str(path))

    doc = DocxExtractor().extract(path)
    assert doc.title == "Anschreiben"  # aus Core-Properties
    assert doc.doc_type == "docx"
    assert "ich bewerbe mich auf die Stelle." in doc.text
    assert "Position | Entwickler" in doc.text


def test_docx_extractor_title_falls_back_to_stem(tmp_path: Path):
    path = tmp_path / "Zeugnis 2024.docx"
    document = Document()
    document.add_paragraph("Sehr gut.")
    document.save(str(path))

    doc = DocxExtractor().extract(path)
    assert doc.title == "Zeugnis 2024"
    assert "Sehr gut." in doc.text


def test_docx_extractor_corrupt_file_does_not_raise(tmp_path: Path):
    path = tmp_path / "kaputt.docx"
    path.write_bytes(b"not really a docx")
    doc = DocxExtractor().extract(path)
    assert doc.doc_type == "docx"
    assert doc.text == ""
    assert doc.title == "kaputt"


def test_pptx_extractor_slides_and_notes(tmp_path: Path):
    path = tmp_path / "Pitch.pptx"
    presentation = Presentation()
    presentation.core_properties.title = "Produktvision"
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Seiton Brain"
    slide.notes_slide.notes_text_frame.text = "Sprechernotiz: Demo zeigen."
    presentation.save(str(path))

    doc = PptxExtractor().extract(path)
    assert doc.title == "Produktvision"  # aus Core-Properties
    assert doc.doc_type == "pptx"
    assert "Seiton Brain" in doc.text
    assert "Sprechernotiz: Demo zeigen." in doc.text


def test_pptx_extractor_title_falls_back_to_stem(tmp_path: Path):
    path = tmp_path / "Quartalsreview.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Zahlen Q1"
    presentation.save(str(path))

    doc = PptxExtractor().extract(path)
    assert doc.title == "Quartalsreview"
    assert "Zahlen Q1" in doc.text


def test_pptx_extractor_corrupt_file_does_not_raise(tmp_path: Path):
    path = tmp_path / "kaputt.pptx"
    path.write_bytes(b"not really a pptx")
    doc = PptxExtractor().extract(path)
    assert doc.doc_type == "pptx"
    assert doc.text == ""
    assert doc.title == "kaputt"
