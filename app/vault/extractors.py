"""Document extraction (E18-1, E18-2, E18-3, E18-5, E18-6).

Engine+adapter pattern for multi-format ingestion: each ``DocumentExtractor``
reads a file group (read-only) and returns plain text for the vault index
(E5-1) and later retrieval/RAG (E17).

Currently tier 1 (directly text-based): Markdown, plain text, PDF (text layer),
plus Office formats Word (.docx) and PowerPoint (.pptx).
OCR (E18-5) and Vision (E18-6) attach as optional adapters (images + PDF scans).
"""

from __future__ import annotations

import logging
from abc import ABC, abstractmethod
from dataclasses import dataclass
from pathlib import Path
from typing import ClassVar

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from app.vault.ocr import IMAGE_OCR_EXTENSIONS, ocr_image, ocr_pdf, ocr_ready, pdf_ocr_ready
from app.vault.reader import _parse_frontmatter
from app.vault.vision import describe_image, format_vision_text, vision_ready

logger = logging.getLogger(__name__)

# Marker for PDFs without an extractable text layer (scans) — OCR hook (E18-5).
PDF_NO_TEXT_TYPE = "pdf_no_text"
PDF_OCR_TYPE = "pdf_ocr"
IMAGE_OCR_TYPE = "image_ocr"
IMAGE_VISION_TYPE = "image_vision"


@dataclass(frozen=True)
class ExtractedDocument:
    """Extraction result: plain text plus metadata."""

    title: str
    text: str
    category: str = ""
    doc_type: str = "text"


class DocumentExtractor(ABC):
    """Adapter interface: read a file and return ``ExtractedDocument``."""

    doc_type: ClassVar[str]
    extensions: ClassVar[tuple[str, ...]]

    @abstractmethod
    def extract(self, path: Path) -> ExtractedDocument:
        """Extract text/metadata from ``path`` (never mutates the file)."""


def _strip_frontmatter(content: str) -> str:
    if not content.startswith("---"):
        return content
    parts = content.split("---", 2)
    return parts[2] if len(parts) > 2 else content


class MarkdownExtractor(DocumentExtractor):
    doc_type = "markdown"
    extensions = (".md", ".markdown")

    def extract(self, path: Path) -> ExtractedDocument:
        content = path.read_text(encoding="utf-8")
        meta = _parse_frontmatter(content)
        title = meta.get("title") or path.stem
        return ExtractedDocument(
            title=title,
            text=_strip_frontmatter(content),
            category=meta.get("category", ""),
            doc_type=self.doc_type,
        )


class PlainTextExtractor(DocumentExtractor):
    doc_type = "text"
    extensions = (".txt", ".text", ".log")

    def extract(self, path: Path) -> ExtractedDocument:
        return ExtractedDocument(
            title=path.stem,
            text=path.read_text(encoding="utf-8"),
            doc_type=self.doc_type,
        )


class PdfExtractor(DocumentExtractor):
    doc_type = "pdf"
    extensions = (".pdf",)

    def extract(self, path: Path) -> ExtractedDocument:
        title = path.stem
        text_parts: list[str] = []
        try:
            reader = PdfReader(str(path))
            meta_title = reader.metadata.title if reader.metadata else None
            if meta_title and meta_title.strip():
                title = meta_title.strip()
            for page in reader.pages:
                extracted = page.extract_text() or ""
                if extracted.strip():
                    text_parts.append(extracted.strip())
        except Exception as exc:  # noqa: BLE001 — broken PDFs must not abort the scan
            logger.warning("PDF-Extraktion fehlgeschlagen fuer %s: %s", path, exc)

        text = "\n\n".join(text_parts).strip()
        if text:
            return ExtractedDocument(title=title, text=text, doc_type=self.doc_type)

        # No text layer (scan) → optional OCR (E18-5), else mark pdf_no_text.
        if pdf_ocr_ready():
            ocr_text = ocr_pdf(path)
            if ocr_text:
                return ExtractedDocument(title=title, text=ocr_text, doc_type=PDF_OCR_TYPE)

        return ExtractedDocument(title=title, text="", doc_type=PDF_NO_TEXT_TYPE)


def _doc_core_title(properties) -> str | None:
    """Title from Office core properties when present."""
    title = getattr(properties, "title", None)
    return title.strip() if title and title.strip() else None


class DocxExtractor(DocumentExtractor):
    """Word documents (.docx). Older binary .doc formats are not supported."""

    doc_type = "docx"
    extensions = (".docx",)

    def extract(self, path: Path) -> ExtractedDocument:
        title = path.stem
        text_parts: list[str] = []
        try:
            document = Document(str(path))
            title = _doc_core_title(document.core_properties) or title
            for paragraph in document.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text.strip())
            # Append tables separately — certificates/invoices often live there.
            for table in document.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        text_parts.append(" | ".join(cells))
        except Exception as exc:  # noqa: BLE001 — broken files must not abort the scan
            logger.warning("DOCX-Extraktion fehlgeschlagen fuer %s: %s", path, exc)

        return ExtractedDocument(
            title=title, text="\n\n".join(text_parts).strip(), doc_type=self.doc_type
        )


class PptxExtractor(DocumentExtractor):
    """PowerPoint presentations (.pptx): slide text + speaker notes."""

    doc_type = "pptx"
    extensions = (".pptx",)

    def extract(self, path: Path) -> ExtractedDocument:
        title = path.stem
        text_parts: list[str] = []
        try:
            presentation = Presentation(str(path))
            title = _doc_core_title(presentation.core_properties) or title
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        text_parts.append(shape.text_frame.text.strip())
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame
                    if notes is not None and notes.text.strip():
                        text_parts.append(notes.text.strip())
        except Exception as exc:  # noqa: BLE001 — broken files must not abort the scan
            logger.warning("PPTX-Extraktion fehlgeschlagen fuer %s: %s", path, exc)

        return ExtractedDocument(
            title=title, text="\n\n".join(text_parts).strip(), doc_type=self.doc_type
        )


class ImageOcrExtractor(DocumentExtractor):
    """Images: OCR (E18-5) when text is found, else Vision LLM (E18-6)."""

    doc_type = IMAGE_OCR_TYPE
    extensions = tuple(sorted(IMAGE_OCR_EXTENSIONS))

    def extract(self, path: Path) -> ExtractedDocument:
        if ocr_ready():
            text = ocr_image(path)
            if text.strip():
                return ExtractedDocument(
                    title=path.stem, text=text, doc_type=IMAGE_OCR_TYPE
                )

        if vision_ready():
            vision = describe_image(path)
            if vision is not None:
                return ExtractedDocument(
                    title=path.stem,
                    text=format_vision_text(vision),
                    doc_type=IMAGE_VISION_TYPE,
                )

        # OCR on but no text / Vision off or failed
        if ocr_ready():
            return ExtractedDocument(title=path.stem, text="", doc_type=IMAGE_OCR_TYPE)
        return ExtractedDocument(title=path.stem, text="", doc_type=IMAGE_VISION_TYPE)


# Alias: kombinierter Bild-Adapter (OCR + Vision).
ImageExtractor = ImageOcrExtractor


# Order decides resolution for ambiguous extensions (unambiguous here).
_EXTRACTORS: tuple[DocumentExtractor, ...] = (
    MarkdownExtractor(),
    PlainTextExtractor(),
    PdfExtractor(),
    DocxExtractor(),
    PptxExtractor(),
)

_IMAGE_EXTRACTOR = ImageOcrExtractor()

# Always-available extensions (without optional OCR/Vision).
SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(
    ext for extractor in _EXTRACTORS for ext in extractor.extensions
)


def image_extraction_ready() -> bool:
    return ocr_ready() or vision_ready()


def get_extractor(path: Path) -> DocumentExtractor | None:
    """Return the matching extractor or ``None`` for unsupported types."""
    suffix = path.suffix.lower()
    for extractor in _EXTRACTORS:
        if suffix in extractor.extensions:
            return extractor
    if suffix in IMAGE_OCR_EXTENSIONS and image_extraction_ready():
        return _IMAGE_EXTRACTOR
    return None


def is_supported(path: Path) -> bool:
    suffix = path.suffix.lower()
    if suffix in SUPPORTED_EXTENSIONS:
        return True
    return suffix in IMAGE_OCR_EXTENSIONS and image_extraction_ready()
